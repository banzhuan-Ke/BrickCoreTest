"""资料库 v2 长文档类解析（docx / pdf / md / txt / doc / ppt / pptx）"""
from __future__ import annotations

import io
import logging
import re
from typing import Any

from app.modules.knowledge.constants import KNOWLEDGE_INGEST_MAX_PROSE_CHARS

logger = logging.getLogger(__name__)

_MUPDF_OFFICE_TYPES = {"doc", "ppt", "pptx"}


def _normalize_section(
    *,
    title: str,
    content: str,
    level: int = 1,
    page: int | None = None,
) -> dict[str, Any]:
    body = (content or "").strip()
    sec: dict[str, Any] = {
        "title": (title or "正文").strip()[:200],
        "level": int(level or 1),
        "content": body,
        "char_count": len(body),
    }
    if page is not None:
        sec["page"] = page
    return sec


def _build_outline(sections: list[dict[str, Any]]) -> list[dict[str, Any]]:
    outline: list[dict[str, Any]] = []
    for sec in sections:
        item: dict[str, Any] = {
            "title": sec.get("title") or "",
            "level": int(sec.get("level") or 1),
            "char_count": int(sec.get("char_count") or 0),
        }
        if sec.get("page") is not None:
            item["page"] = sec["page"]
        outline.append(item)
    return outline


def _sections_to_fulltext(sections: list[dict[str, Any]]) -> str:
    parts: list[str] = []
    for sec in sections:
        title = (sec.get("title") or "").strip()
        body = (sec.get("content") or "").strip()
        if title and body:
            parts.append(f"## {title}\n{body}")
        elif body:
            parts.append(body)
        elif title:
            parts.append(f"## {title}")
    return "\n\n".join(parts).strip()


def _apply_prose_size_guard(
    sections: list[dict[str, Any]],
    *,
    max_chars: int,
    warnings: list[str],
) -> tuple[list[dict[str, Any]], str]:
    fulltext = _sections_to_fulltext(sections)
    if len(fulltext) <= max_chars:
        return sections, fulltext

    warnings.append(
        f"文档正文超过上限 {max_chars} 字符，已截断；完整内容请查看源文件"
    )
    kept: list[dict[str, Any]] = []
    used = 0
    budget_exhausted = False
    for sec in sections:
        title = (sec.get("title") or "").strip()
        body = (sec.get("content") or "").strip()
        prefix = f"## {title}\n" if title else ""
        piece = f"{prefix}{body}".strip()
        if budget_exhausted:
            kept.append({**sec, "content": "", "char_count": 0})
            continue
        if not piece:
            kept.append({**sec, "content": "", "char_count": 0})
            continue
        remaining = max_chars - used
        if remaining <= 0:
            budget_exhausted = True
            kept.append({**sec, "content": "", "char_count": 0})
            continue
        if len(piece) <= remaining:
            kept.append(sec)
            used += len(piece) + 2
        else:
            truncated_body = body[: max(0, remaining - len(prefix))]
            kept.append(
                {
                    **sec,
                    "content": truncated_body,
                    "char_count": len(truncated_body),
                }
            )
            used = max_chars
            budget_exhausted = True
    return kept, _sections_to_fulltext(kept)[:max_chars]


def _sections_from_loader(result: Any) -> list[dict[str, Any]]:
    from app.modules.ai.requirement_document import section_text

    blocks = result.blocks or []
    raw_sections = result.sections or []
    sections: list[dict[str, Any]] = []

    if raw_sections and blocks:
        for sec in raw_sections:
            if not isinstance(sec, dict):
                continue
            title = (sec.get("title") or "正文").strip()
            level = int(sec.get("level") or 1)
            page = sec.get("page")
            if page is None:
                m = re.match(r"第\s*(\d+)\s*页", title)
                if m:
                    page = int(m.group(1))
            body = section_text(blocks, sec)
            sections.append(_normalize_section(title=title, content=body, level=level, page=page))
    elif result.text:
        sections.append(_normalize_section(title="全文", content=result.text, level=1))

    return sections


def _parse_pptx_sections(content: bytes) -> list[dict[str, Any]]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    sections: list[dict[str, Any]] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    texts.append(t)
        if not texts:
            continue
        title = texts[0].split("\n")[0].strip()[:200] or f"Slide {slide_idx}"
        body = "\n".join(texts)
        sections.append(_normalize_section(title=title, content=body, level=1, page=slide_idx))
    return sections


def _parse_mupdf_office_sections(content: bytes, file_name: str, filetype: str) -> list[dict[str, Any]]:
    import fitz

    doc = fitz.open(stream=content, filetype=filetype, filename=file_name or f"file.{filetype}")
    sections: list[dict[str, Any]] = []
    try:
        for page_idx in range(len(doc)):
            page_text = (doc[page_idx].get_text() or "").strip()
            if not page_text:
                continue
            sections.append(
                _normalize_section(
                    title=f"第 {page_idx + 1} 页",
                    content=page_text,
                    level=1,
                    page=page_idx + 1,
                )
            )
    finally:
        doc.close()
    return sections


def parse_prose_v2(content: bytes, file_name: str, ext: str) -> dict[str, Any]:
    """解析长文档为 v2 结构（sections + fulltext + outline）。"""
    warnings: list[str] = []
    sections: list[dict[str, Any]] = []
    page_count = 0
    fmt = ext

    loader_meta: dict[str, Any] | None = None
    if ext in ("docx", "pdf", "md", "txt"):
        from app.modules.ai.document_loader import load_document
        from app.modules.knowledge.knowledge_image_parse import prepare_loader_meta_from_result

        result = load_document(content, file_name)
        fmt = result.source_type or ext
        warnings.extend(result.warnings or [])
        page_count = int(result.page_count or 0)
        sections = _sections_from_loader(result)
        loader_meta = prepare_loader_meta_from_result(result)
    elif ext == "pptx":
        try:
            sections = _parse_pptx_sections(content)
        except Exception as ex:
            logger.warning("[knowledge_ingest_prose] pptx shapes failed: %s", ex)
            sections = _parse_mupdf_office_sections(content, file_name, "pptx")
            fmt = "pptx"
    elif ext in _MUPDF_OFFICE_TYPES:
        sections = _parse_mupdf_office_sections(content, file_name, ext)
        page_count = len(sections)
    else:
        raise ValueError(f"不支持的 prose 格式: {ext}")

    if not sections:
        sections.append(_normalize_section(title="全文", content="", level=1))

    sections, fulltext = _apply_prose_size_guard(
        sections,
        max_chars=KNOWLEDGE_INGEST_MAX_PROSE_CHARS,
        warnings=warnings,
    )
    outline = _build_outline(sections)
    char_count = len(fulltext)

    payload: dict[str, Any] = {
        "format": fmt,
        "sections": sections,
        "fulltext": fulltext,
        "char_count": char_count,
        "parse_warnings": warnings,
        "chunk_count": max(len(sections), 1 if fulltext else 0),
        "structured": {
            "kind": "prose",
            "outline": outline,
        },
    }
    if page_count > 0:
        payload["page_count"] = page_count
    elif sections and all(sec.get("page") for sec in sections):
        payload["page_count"] = len(sections)
    if loader_meta:
        payload["loader_meta"] = loader_meta
    return payload


def build_prose_sections_json_v2(
    *,
    sections: list[dict[str, Any]],
    fmt: str,
    parse_warnings: list[str] | None = None,
) -> dict[str, Any]:
    """由已有 sections 构建 prose v2 sections_json（归档/需求同步用）。"""
    from app.modules.knowledge.constants import SECTIONS_JSON_VERSION_V2

    warnings = list(parse_warnings or [])
    normalized: list[dict[str, Any]] = []
    for raw in sections:
        if not isinstance(raw, dict):
            continue
        body = (raw.get("content") or raw.get("text") or raw.get("_text_buf") or "").strip()
        normalized.append(
            _normalize_section(
                title=(raw.get("title") or "正文"),
                content=body,
                level=int(raw.get("level") or 1),
                page=raw.get("page"),
            )
        )
    if not normalized:
        normalized.append(_normalize_section(title="全文", content="", level=1))

    normalized, fulltext = _apply_prose_size_guard(
        normalized,
        max_chars=KNOWLEDGE_INGEST_MAX_PROSE_CHARS,
        warnings=warnings,
    )
    outline = _build_outline(normalized)
    return {
        "version": SECTIONS_JSON_VERSION_V2,
        "format": fmt,
        "char_count": len(fulltext),
        "parse_warnings": warnings,
        "sections": normalized,
        "fulltext": fulltext,
        "structured": {
            "kind": "prose",
            "outline": outline,
        },
    }
